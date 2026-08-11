#!/usr/bin/env python3
"""Text-track extractor for the powerpoint-to-md skill.

Reads a .pptx and produces an intermediate ``extract.json`` in the given tmp
directory, plus dumps every embedded image to ``images-out``.

Sources of truth:
  * Microsoft MarkItDown  — deck-level structured markdown (headings, bullets, tables)
  * python-pptx           — per-slide notes, tables, embedded images, chart series,
                            layout names, native chart data

MarkItDown produces a single markdown blob for the whole deck. We split that
blob per-slide using its own ``<!-- Slide number: N -->`` markers so each slide
in ``extract.json`` gets its own ``content_markdown`` field.

Everything downstream (the calling agent) reads ``extract.json``.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
import traceback
from dataclasses import dataclass, field, asdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any


def log(msg: str) -> None:
    print(f"[extract_text] {msg}", file=sys.stderr)


# ---------- data model ----------


@dataclass
class ChartInfo:
    type: str
    title: str
    categories: list[str] = field(default_factory=list)
    series: list[dict[str, Any]] = field(default_factory=list)
    note: str = ""


@dataclass
class TableInfo:
    headers: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)


@dataclass
class ImageInfo:
    path: str
    alt: str
    shape_name: str


@dataclass
class SlideInfo:
    index: int
    layout: str = ""
    title: str = ""
    subtitle: str = ""
    content_markdown: str = ""
    bullets: list[str] = field(default_factory=list)
    tables: list[TableInfo] = field(default_factory=list)
    charts: list[ChartInfo] = field(default_factory=list)
    images: list[ImageInfo] = field(default_factory=list)
    notes: str = ""
    image_render: str = ""  # filled in by render_slides.sh — path relative to output


@dataclass
class DeckInfo:
    source: str
    title: str
    author: str
    slide_count: int
    generated_at: str


# ---------- markitdown split ----------


_MD_SLIDE_MARKER = re.compile(r"<!--\s*Slide number:\s*(\d+)\s*-->", re.IGNORECASE)


def split_markitdown_per_slide(md_text: str) -> dict[int, str]:
    """Split a MarkItDown pptx output into per-slide chunks.

    MarkItDown emits an HTML comment ``<!-- Slide number: N -->`` before each
    slide's content. We split on that marker. If markers are absent, we return
    a single mapping ``{1: md_text}``.
    """
    if not md_text:
        return {}

    matches = list(_MD_SLIDE_MARKER.finditer(md_text))
    if not matches:
        return {1: md_text.strip()}

    per_slide: dict[int, str] = {}
    for i, m in enumerate(matches):
        slide_num = int(m.group(1))
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(md_text)
        chunk = md_text[start:end].strip()
        per_slide[slide_num] = chunk
    return per_slide


# ---------- python-pptx extraction ----------


def _text_from_shape(shape) -> str:
    """Return the full text of any shape that has a text frame."""
    if not getattr(shape, "has_text_frame", False):
        return ""
    tf = shape.text_frame
    parts: list[str] = []
    for p in tf.paragraphs:
        line = "".join(r.text or "" for r in p.runs)
        if not line:
            line = p.text or ""
        parts.append(line)
    return "\n".join(parts).strip()


def _iter_all_shapes(shapes):
    """Recurse into group shapes to yield every leaf shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # local import to avoid hard dep at module load

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape.shapes)
        else:
            yield shape


def _bullets_from_shape(shape) -> list[str]:
    if not getattr(shape, "has_text_frame", False):
        return []
    out: list[str] = []
    for p in shape.text_frame.paragraphs:
        text = "".join(r.text or "" for r in p.runs) or (p.text or "")
        text = text.strip()
        if text:
            out.append(text)
    return out


def _extract_table(shape) -> TableInfo | None:
    if not getattr(shape, "has_table", False):
        return None
    tbl = shape.table
    rows = list(tbl.rows)
    if not rows:
        return None
    headers: list[str] = []
    body: list[list[str]] = []
    for r_idx, row in enumerate(rows):
        cells = [(cell.text or "").strip().replace("\n", " ") for cell in row.cells]
        if r_idx == 0:
            headers = cells
        else:
            body.append(cells)
    return TableInfo(headers=headers, rows=body)


def _extract_chart(shape) -> ChartInfo | None:
    if not getattr(shape, "has_chart", False):
        return None
    try:
        chart = shape.chart
    except Exception:  # noqa: BLE001
        return None

    chart_type = ""
    try:
        if chart.chart_type is not None:
            # python-pptx enum repr is "XL_CHART_TYPE.BAR_CLUSTERED (57)"; take just the name.
            raw = str(chart.chart_type)
            raw = raw.split(".", 1)[-1]  # drop "XL_CHART_TYPE." prefix
            raw = raw.split(" ", 1)[0]   # drop " (57)" suffix
            chart_type = raw
    except Exception:  # noqa: BLE001
        chart_type = ""

    title = ""
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text.strip()
    except Exception:  # noqa: BLE001
        title = ""

    info = ChartInfo(type=chart_type, title=title)

    # Categories — from the first plot.
    try:
        plots = list(chart.plots)
        if plots:
            info.categories = [str(c) for c in plots[0].categories]
    except Exception:  # noqa: BLE001
        pass

    # Series — name + values.
    try:
        for s in chart.series:
            name = getattr(s, "name", "") or ""
            values: list[Any] = []
            try:
                values = [v for v in s.values]
            except Exception:  # noqa: BLE001
                values = []
            info.series.append({"name": name, "values": values})
    except Exception as exc:  # noqa: BLE001
        info.note = f"chart series extraction partial: {exc}"

    return info


def _extract_images(shape, slide_index: int, out_dir: Path, seen: set[str]) -> list[ImageInfo]:
    """Dump embedded images from a shape (recursively for groups) to ``out_dir``."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    results: list[ImageInfo] = []
    shapes_iter = [shape]
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        shapes_iter = list(_iter_all_shapes([shape]))

    for s in shapes_iter:
        picture = None
        if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            picture = s
        elif hasattr(s, "image"):
            try:
                _ = s.image
                picture = s
            except Exception:  # noqa: BLE001
                picture = None
        if picture is None:
            continue
        try:
            image = picture.image
        except Exception:  # noqa: BLE001
            continue

        blob = image.blob
        ext = (image.ext or "png").lower().lstrip(".")
        digest = hashlib.sha1(blob).hexdigest()[:8]  # noqa: S324 (not security)
        pic_index = len(results) + 1
        fname = f"slide-{slide_index:02d}-pic{pic_index}-{digest}.{ext}"
        dest = out_dir / fname
        if fname not in seen:
            dest.write_bytes(blob)
            seen.add(fname)
        alt = ""
        try:
            alt = (picture.name or "").strip()
        except Exception:  # noqa: BLE001
            pass
        shape_name = getattr(picture, "name", "") or ""
        results.append(ImageInfo(path=str(dest), alt=alt, shape_name=shape_name))
    return results


def extract_with_pptx(pptx_path: Path, images_out: Path) -> tuple[DeckInfo, list[SlideInfo]]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    cp = prs.core_properties

    deck = DeckInfo(
        source=str(pptx_path),
        title=(cp.title or "").strip(),
        author=(cp.author or "").strip(),
        slide_count=len(prs.slides),
        generated_at=datetime.now(timezone.utc).replace(microsecond=0).isoformat(),
    )

    images_out.mkdir(parents=True, exist_ok=True)
    seen: set[str] = set()

    slides: list[SlideInfo] = []
    for idx, slide in enumerate(prs.slides, start=1):
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name or ""
        except Exception:  # noqa: BLE001
            layout_name = ""

        title = ""
        subtitle = ""
        try:
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
        except Exception:  # noqa: BLE001
            title = ""

        bullets: list[str] = []
        tables: list[TableInfo] = []
        charts: list[ChartInfo] = []
        images: list[ImageInfo] = []

        for shape in _iter_all_shapes(slide.shapes):
            # Subtitle: placeholder idx 1 in Title-Slide layouts.
            if getattr(shape, "is_placeholder", False):
                try:
                    ph = shape.placeholder_format
                    if ph is not None and ph.idx == 1 and not subtitle:
                        sub_text = _text_from_shape(shape)
                        if sub_text and sub_text != title:
                            subtitle = sub_text
                except Exception:  # noqa: BLE001
                    pass

            # Bullets — text of any text-bearing shape that isn't the title.
            if getattr(shape, "has_text_frame", False):
                if title and _text_from_shape(shape) == title:
                    pass
                else:
                    for b in _bullets_from_shape(shape):
                        if b == title:
                            continue
                        if b == subtitle:
                            continue
                        bullets.append(b)

            tbl = _extract_table(shape)
            if tbl is not None:
                tables.append(tbl)

            chart = _extract_chart(shape)
            if chart is not None:
                charts.append(chart)

            images.extend(_extract_images(shape, idx, images_out, seen))

        # Speaker notes.
        notes = ""
        try:
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf is not None:
                    notes = (notes_tf.text or "").strip()
        except Exception:  # noqa: BLE001
            notes = ""

        slides.append(
            SlideInfo(
                index=idx,
                layout=layout_name,
                title=title,
                subtitle=subtitle,
                bullets=_dedup_preserve_order(bullets),
                tables=tables,
                charts=charts,
                images=images,
                notes=notes,
            )
        )
    return deck, slides


def _dedup_preserve_order(items: list[str]) -> list[str]:
    seen: set[str] = set()
    out: list[str] = []
    for it in items:
        if it in seen:
            continue
        seen.add(it)
        out.append(it)
    return out


# ---------- markitdown ----------


def extract_with_markitdown(pptx_path: Path) -> str:
    try:
        from markitdown import MarkItDown  # type: ignore
    except Exception as exc:  # noqa: BLE001
        log(f"WARN: markitdown import failed: {exc} — proceeding with python-pptx only")
        return ""

    try:
        md = MarkItDown(enable_plugins=False)
        result = md.convert(str(pptx_path))
        # Newer versions use .text_content; older/alt: .markdown; guard both.
        text = getattr(result, "text_content", None) or getattr(result, "markdown", None) or ""
        return text
    except Exception as exc:  # noqa: BLE001
        log(f"WARN: markitdown conversion failed: {exc} — proceeding with python-pptx only")
        traceback.print_exc(file=sys.stderr)
        return ""


# ---------- output ----------


def build_output(
    deck: DeckInfo,
    slides: list[SlideInfo],
    md_by_slide: dict[int, str],
) -> dict[str, Any]:
    slide_dicts: list[dict[str, Any]] = []
    for s in slides:
        s.content_markdown = md_by_slide.get(s.index, "")
        d = asdict(s)
        slide_dicts.append(d)

    return {
        "deck": asdict(deck),
        "slides": slide_dicts,
    }


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Text-track extractor for pptx → extract.json")
    ap.add_argument("--input", required=True, help="Path to input .pptx")
    ap.add_argument("--tmp-dir", required=True, help="Tmp dir to write extract.json into")
    ap.add_argument("--images-out", required=True, help="Directory to dump embedded images")
    args = ap.parse_args(argv)

    input_path = Path(args.input).resolve()
    tmp_dir = Path(args.tmp_dir).resolve()
    images_out = Path(args.images_out).resolve()

    if not input_path.exists():
        log(f"ERROR: input not found: {input_path}")
        return 2

    tmp_dir.mkdir(parents=True, exist_ok=True)
    images_out.mkdir(parents=True, exist_ok=True)

    log(f"input: {input_path}")
    log(f"tmp:   {tmp_dir}")
    log(f"images:{images_out}")

    # Track A1: python-pptx (authoritative for notes, tables, charts, images).
    try:
        deck, slides = extract_with_pptx(input_path, images_out)
    except Exception as exc:  # noqa: BLE001
        log(f"ERROR: python-pptx extraction failed: {exc}")
        traceback.print_exc(file=sys.stderr)
        return 3

    log(f"parsed {deck.slide_count} slides via python-pptx")

    # Track A2: markitdown (structured per-slide markdown).
    md_text = extract_with_markitdown(input_path)
    md_by_slide = split_markitdown_per_slide(md_text)
    log(f"markitdown produced markdown for {len(md_by_slide)} slide chunk(s)")

    # Compose.
    payload = build_output(deck, slides, md_by_slide)

    out_json = tmp_dir / "extract.json"
    out_json.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2, default=str),
        encoding="utf-8",
    )
    log(f"wrote {out_json}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
